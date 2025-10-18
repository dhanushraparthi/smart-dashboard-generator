# app.py
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os

# Optional OpenAI
try:
    import openai
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

# ---------------- Page Setup ----------------
st.set_page_config(page_title="Smart Dashboard Generator", layout="wide")
st.title("📊 Smart Dashboard Generator")
st.markdown("Upload CSV, Excel, or JSON datasets to generate KPIs, dashboards, AI insights, and PDF/PPTX reports.")

# ---------------- Sidebar ----------------
st.sidebar.header("Upload & Filters")
uploaded_file = st.sidebar.file_uploader("Upload dataset", type=["csv", "xlsx", "xls", "json"])

chart_types = st.sidebar.multiselect(
    "Select chart types", ["Bar Chart", "Line Chart", "Scatter Chart"], default=["Bar Chart"]
)

# ---------------- Load Dataset ----------------
def load_file(uploaded):
    name = uploaded.name.lower()
    if name.endswith(".csv"):
        return pd.read_csv(uploaded)
    elif name.endswith((".xlsx", ".xls")):
        return pd.read_excel(uploaded)
    elif name.endswith(".json"):
        return pd.read_json(uploaded)
    else:
        raise ValueError("Unsupported file type")

if not uploaded_file:
    st.info("Please upload a dataset to begin.")
    st.stop()

try:
    df = load_file(uploaded_file)
except Exception as e:
    st.error(f"Failed to load dataset: {e}")
    st.stop()

df.columns = df.columns.str.strip().str.lower().str.replace(" ", "_")
st.success("Dataset loaded successfully!")
st.dataframe(df.head())

# ---------------- Sidebar Filters ----------------
filter_cols = st.sidebar.multiselect("Filter by Columns", df.columns.tolist())
filtered_df = df.copy()

for col in filter_cols:
    unique_vals = df[col].dropna().unique().tolist()
    selected_vals = st.sidebar.multiselect(f"Select {col}", unique_vals, default=unique_vals)
    filtered_df = filtered_df[filtered_df[col].isin(selected_vals)]

# ---------------- KPIs ----------------
st.subheader("Key Metrics")
kpi_cols = st.columns(4)
total_sales = filtered_df["sales"].sum() if "sales" in filtered_df.columns else None
total_profit = filtered_df["profit"].sum() if "profit" in filtered_df.columns else None
avg_discount = filtered_df["discount"].mean() if "discount" in filtered_df.columns else None
total_quantity = filtered_df["quantity"].sum() if "quantity" in filtered_df.columns else None

if total_sales is not None:
    kpi_cols[0].metric("Total Sales", f"${total_sales:,.2f}")
if total_profit is not None:
    kpi_cols[1].metric("Total Profit", f"${total_profit:,.2f}")
if avg_discount is not None:
    kpi_cols[2].metric("Avg Discount", f"{avg_discount:.2%}")
if total_quantity is not None:
    kpi_cols[3].metric("Total Quantity", f"{int(total_quantity):,}")

# ---------------- Dashboards ----------------
st.subheader("Visual Dashboards")
plots = {}
numeric_cols = filtered_df.select_dtypes(include=np.number).columns.tolist()
categorical_cols = filtered_df.select_dtypes(include="object").columns.tolist()

# Bar charts for categorical columns
for col in categorical_cols:
    if "Bar Chart" in chart_types:
        vc = filtered_df[col].value_counts().reset_index()
        vc.columns = [col, "count"]  # Fix for Plotly
        fig = px.bar(vc, x=col, y="count", text_auto=True, title=f"{col.capitalize()} Distribution", template="plotly_dark")
        st.plotly_chart(fig, use_container_width=True)
        plots[col] = fig

# Line charts for numeric columns
for col in numeric_cols:
    if "Line Chart" in chart_types:
        fig = px.line(filtered_df, y=col, title=f"{col.capitalize()} Trend", template="plotly_dark")
        st.plotly_chart(fig, use_container_width=True)
        plots[col] = fig

# Scatter chart for first two numeric columns
if len(numeric_cols) >= 2 and "Scatter Chart" in chart_types:
    x_col, y_col = numeric_cols[:2]
    fig = px.scatter(filtered_df, x=x_col, y=y_col, color=categorical_cols[0] if categorical_cols else None,
                     title=f"{y_col} vs {x_col}", template="plotly_dark")
    st.plotly_chart(fig, use_container_width=True)
    plots[f"{x_col}_vs_{y_col}"] = fig

# ---------------- AI Insights ----------------
st.subheader("AI Insights")

def heuristic_insights(df):
    insights = []
    if "region" in df.columns and "sales" in df.columns:
        top_region = df.groupby("region")["sales"].sum().idxmax()
        insights.append(f"Top region by sales: {top_region}")
    if "category" in df.columns and "profit" in df.columns:
        top_category = df.groupby("category")["profit"].sum().idxmax()
        insights.append(f"Most profitable category: {top_category}")
    if "discount" in df.columns and "profit" in df.columns:
        corr = df["discount"].corr(df["profit"])
        insights.append(f"Discount vs Profit correlation: {corr:.2f}")
    if "sales" in df.columns and "profit" in df.columns:
        margin = df["profit"].sum() / df["sales"].sum() if df["sales"].sum() != 0 else 0
        insights.append(f"Overall profit margin: {margin:.2%}")
    if not insights:
        insights.append("Not enough data for detailed insights.")
    return insights

ai_key = os.environ.get("OPENAI_API_KEY")
ai_text = ""
if OPENAI_AVAILABLE and ai_key:
    try:
        openai.api_key = ai_key
        prompt = "Summarize dataset insights in 3 bullet points:\n" + "\n".join(heuristic_insights(filtered_df))
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=200,
            temperature=0.2
        )
        ai_text = response.choices[0].message.content.strip()
    except Exception:
        ai_text = "\n".join(heuristic_insights(filtered_df))
else:
    ai_text = "\n".join(heuristic_insights(filtered_df))

for line in ai_text.split("\n"):
    st.markdown(f"- {line.strip()}")

# ---------------- PDF Export ----------------
st.subheader("Export Reports")

def create_pdf(kpis, plots, insights):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, "Smart Dashboard Report", ln=True, align="C")
    pdf.ln(5)
    pdf.set_font("Arial", size=12)
    for kpi in kpis:
        pdf.multi_cell(0, 6, kpi)
    pdf.ln(5)
    pdf.set_font("Arial", "B", 12)
    pdf.cell(0, 6, "Insights:", ln=True)
    pdf.set_font("Arial", size=11)
    for line in insights.split("\n"):
        pdf.multi_cell(0, 6, line)
    pdf.ln(5)
    for name, fig in plots.items():
        img_bytes = fig.to_image(format="png", engine="kaleido")
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        pdf.image(tmp.name, w=170)
        os.remove(tmp.name)
    return pdf.output(dest="S").encode("latin-1", errors="ignore")

kpis_list = []
if total_sales: kpis_list.append(f"Total Sales: ${total_sales:,.2f}")
if total_profit: kpis_list.append(f"Total Profit: ${total_profit:,.2f}")
if avg_discount: kpis_list.append(f"Avg Discount: {avg_discount:.2%}")
if total_quantity: kpis_list.append(f"Total Quantity: {int(total_quantity):,}")

if st.button("Download PDF Report"):
    pdf_bytes = create_pdf(kpis_list, plots, ai_text)
    st.download_button("Download PDF", data=pdf_bytes, file_name="dashboard_report.pdf", mime="application/pdf")

# ---------------- PPTX Export ----------------
def create_pptx(kpis, plots, insights):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Smart Dashboard Report"

    top = Inches(1)
    left = Inches(0.5)
    width = Inches(9)

    for kpi in kpis:
        textbox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        textbox.text = kpi
        top += Inches(0.3)

    top += Inches(0.2)
    textbox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    textbox.text = "Insights:"
    top += Inches(0.3)
    for line in insights.split("\n"):
        textbox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        textbox.text = line
        top += Inches(0.3)

    for name, fig in plots.items():
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = name
        img_bytes = fig.to_image(format="png", engine="kaleido")
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        slide.shapes.add_picture(tmp.name, Inches(0.5), Inches(1.5), width=Inches(9))
        os.remove(tmp.name)

    buf = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(buf.name)
    buf.seek(0)
    data = buf.read()
    buf.close()
    os.remove(buf.name)
    return data

if st.button("Download PPTX Report"):
    pptx_data = create_pptx(kpis_list, plots, ai_text)
    st.download_button("Download PPTX", data=pptx_data, file_name="dashboard_report.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
